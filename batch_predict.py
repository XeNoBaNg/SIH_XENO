# batch_predict.py
import os
import argparse
import csv
from PIL import Image
from tqdm import tqdm
from pptx import Presentation
import pptx.util
import io
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

# Import your CivicIssueFuser
from models.ensemble.fuser import CivicIssueFuser

def generate_report(image_folder: str, weights_path: str, num_classes: int, output_file: str, file_format: str, metadata_csv: str = None):
    """
    Processes a batch of images and generates a report with a custom slide layout.
    """
    print(f"--- Starting Batch Prediction (Format: {file_format.upper()}) ---")
    
    fuser = CivicIssueFuser(
        backbone_weights=weights_path,
        num_classes=num_classes,
    )
    
    descriptions = {}
    if metadata_csv:
        print(f"Loading descriptions from {metadata_csv}...")
        try:
            with open(metadata_csv, 'r', encoding='utf-8-sig') as f:
                reader = csv.DictReader(f)
                for row in reader:
                    descriptions[row['filename']] = row['description']
        except FileNotFoundError:
            print(f"Warning: Metadata file not found at {metadata_csv}. Proceeding without descriptions.")

    image_extensions = ('.jpg', '.jpeg', '.png', '.webp')
    image_files = [f for f in os.listdir(image_folder) if f.lower().endswith(image_extensions)]
    
    if not image_files:
        print(f"Error: No images found in the folder: {image_folder}")
        return
        
    print(f"Found {len(image_files)} images to process.")
    
    if file_format == 'pptx':
        prs = Presentation()
        # --- LAYOUT CHANGE: Use "Title Only" layout for more control ---
        slide_layout = prs.slide_layouts[5]
    else: # CSV
        csv_file = open(output_file, 'w', newline='', encoding='utf-8')
        fieldnames = ['filename', 'user_description', 'prediction', 'confidence', 'ai_caption', 'is_relevant', 'relevance_score']
        writer = csv.DictWriter(csv_file, fieldnames=fieldnames)
        writer.writeheader()

    for filename in tqdm(image_files, desc="Processing images"):
        try:
            image_path = os.path.join(image_folder, filename)
            image = Image.open(image_path).convert("RGB")
            user_desc = descriptions.get(filename, "N/A")
            result = fuser.predict(image, user_description=user_desc)
            
            if file_format == 'pptx':
                slide = prs.slides.add_slide(slide_layout)
                
                # Set Title (this placeholder exists in the "Title Only" layout)
                title = slide.shapes.title
                title.text = f"Analysis for: {filename}"
                
                # Convert image to PNG in memory
                image_stream = io.BytesIO()
                image.save(image_stream, format="PNG")
                image_stream.seek(0)
                
                # --- LAYOUT CHANGE: Manually add and position the image ---
                # Position arguments are (left, top, width)
                # Slide is typically 10" wide by 7.5" high
                img_width = pptx.util.Inches(4.5)
                img_left = (prs.slide_width - img_width) / 2 # Center horizontally
                slide.shapes.add_picture(image_stream, img_left, pptx.util.Inches(1.5), width=img_width)

                # --- LAYOUT CHANGE: Manually add and position the text box ---
                txBox_left = pptx.util.Inches(1.0)
                txBox_top = pptx.util.Inches(4.5)
                txBox_width = pptx.util.Inches(8.0)
                txBox_height = pptx.util.Inches(2.5)
                text_box = slide.shapes.add_textbox(txBox_left, txBox_top, txBox_width, txBox_height)
                
                tf = text_box.text_frame
                tf.clear()
                
                # Apply formatting
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.word_wrap = True

                relevance = result.get('relevance', {})
                spam_status = "Not Spam" if relevance.get('is_relevant') else "Potential Spam"
                relevance_score = relevance.get('relevance_score', 0)
                
                detections_list = result.get('detections', [])
                detected_items = [f"{d['label']} ({d['score']:.2f})" for d in detections_list]
                detections_str = ", ".join(detected_items) if detected_items else "None"
                
                content = (
                    f"User Description: {user_desc}\n\n"
                    f"Spam Check: {spam_status} (Score: {relevance_score:.2f})\n"
                    f"------------------------------------\n"
                    f"AI Model Prediction: {result.get('prediction', 'error')}\n"
                    f"Confidence: {result.get('confidence', 0):.2%}\n"
                    f"AI Generated Caption: {result.get('ai_caption', '')}\n"
                    f"Detected Objects: {detections_str}"
                )
                
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = pptx.util.Pt(12)
                p.alignment = PP_ALIGN.CENTER
            
            else: # CSV Logic
                pass

        except Exception as e:
            print(f"\nError processing {filename}: {e}")

    if file_format == 'pptx':
        prs.save(output_file)
    else:
        csv_file.close()

    print(f"\n--- Batch Prediction Complete ---")
    print(f"✅ Report saved successfully to: {output_file}")


if __name__ == '__main__':
    # (Argument parsing code remains the same)
    parser = argparse.ArgumentParser(description="Run batch predictions on a folder of images.")
    parser.add_argument('--image_folder', type=str, required=True, help="Path to the folder containing images.")
    parser.add_argument('--weights', type=str, required=True, help="Path to the trained model weights file.")
    parser.add_argument('--num_classes', type=int, required=True, help="The number of classes the model was trained on.")
    parser.add_argument('--output_file', type=str, required=True, help="Path to save the output report file.")
    parser.add_argument('--format', type=str, choices=['pptx', 'csv'], default='csv', help="Format of the output report.")
    parser.add_argument('--metadata_csv', type=str, required=False, help="Optional path to a CSV file with descriptions.")
    
    args = parser.parse_args()
    
    generate_report(args.image_folder, args.weights, args.num_classes, args.output_file, args.format, args.metadata_csv)